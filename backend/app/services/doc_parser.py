from __future__ import annotations

import io
import re
from dataclasses import dataclass, field

from app.core.config import settings


@dataclass
class ParseResult:
    text: str
    pages: list[str] = field(default_factory=list)
    used_ocr: bool = False
    error: str | None = None


def _clean(text: str) -> str:
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _strip_markdown_front_matter(text: str) -> str:
    if not text.startswith("---"):
        return text
    match = re.match(r"^---\s*\r?\n.*?\r?\n---\s*(?:\r?\n|$)", text, flags=re.DOTALL)
    return text[match.end() :] if match else text


def _inline_markdown_text(token) -> str:
    parts: list[str] = []
    for child in token.children or []:
        if child.type in {"image", "html_inline"}:
            continue
        if child.type in {"softbreak", "hardbreak"}:
            parts.append(" ")
        elif child.type in {"text", "code_inline"}:
            parts.append(child.content)
    return re.sub(r"\s+", " ", "".join(parts)).strip()


def parse_markdown(data: bytes) -> ParseResult:
    try:
        source = data.decode("utf-8-sig")
    except UnicodeDecodeError:
        return ParseResult(text="", error="Markdown 文件必须使用 UTF-8 编码")

    from markdown_it import MarkdownIt

    source = _strip_markdown_front_matter(source)
    source = re.sub(r"<!--.*?-->", "", source, flags=re.DOTALL)
    tokens = MarkdownIt("commonmark", {"html": False}).enable("table").parse(source)
    blocks: list[str] = []
    heading_path: list[str] = []
    table_row: list[str] = []
    in_table = False
    in_quote = False
    list_depth = 0
    heading_level: int | None = None

    for token in tokens:
        if token.type in {"fence", "code_block", "html_block"}:
            continue
        if token.type == "heading_open":
            heading_level = int(token.tag[1:])
            continue
        if token.type == "heading_close":
            heading_level = None
            continue
        if token.type == "table_open":
            in_table = True
            blocks.append("[表格]")
            continue
        if token.type == "table_close":
            in_table = False
            continue
        if token.type == "tr_open":
            table_row = []
            continue
        if token.type == "tr_close":
            if table_row:
                blocks.append(" | ".join(table_row))
            table_row = []
            continue
        if token.type == "blockquote_open":
            in_quote = True
            continue
        if token.type == "blockquote_close":
            in_quote = False
            continue
        if token.type in {"bullet_list_open", "ordered_list_open"}:
            list_depth += 1
            continue
        if token.type in {"bullet_list_close", "ordered_list_close"}:
            list_depth = max(list_depth - 1, 0)
            continue
        if token.type != "inline":
            continue

        text = _inline_markdown_text(token)
        if not text:
            continue
        if heading_level is not None:
            heading_path = heading_path[: heading_level - 1] + [text]
            continue
        if in_table:
            table_row.append(text)
            continue
        prefix = ""
        if heading_path:
            prefix += f"【{' > '.join(heading_path)}】 "
        if in_quote:
            prefix += "引用："
        if list_depth:
            prefix += "• "
        blocks.append(prefix + text)

    text = "\n\n".join(blocks).strip()
    if not text:
        return ParseResult(text="", error="Markdown 文档中没有可用于出题的文本")
    return ParseResult(text=text, pages=[text])


def parse_pdf(data: bytes) -> ParseResult:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text("text") or "")
    text = _clean("\n".join(pages))
    used_ocr = False
    if settings.enable_ocr and len(text) < 80:
        ocr_pages = []
        try:
            from app.services.ocr import ocr_pixmap

            for i, page in enumerate(doc):
                if i >= 10:
                    break
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                ocr_pages.append(ocr_pixmap(pix.tobytes("png")))
            ocr_text = _clean("\n".join(ocr_pages))
            if len(ocr_text) > len(text):
                text, pages, used_ocr = ocr_text, ocr_pages, True
        except Exception:
            pass
    if not text.strip():
        return ParseResult(text="", pages=pages, used_ocr=used_ocr, error="未能提取到可选中文本，扫描件请检查或稍后使用 OCR")
    tables = []
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables() or []:
                    rows = [" | ".join(cell or "" for cell in row) for row in table]
                    tables.append("\n".join(rows))
    except Exception:
        pass
    if tables:
        text += "\n\n[表格]\n" + "\n\n".join(tables)
    return ParseResult(text=text, pages=pages, used_ocr=used_ocr)


def parse_docx(data: bytes) -> ParseResult:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text = _clean("\n".join(parts))
    if not text:
        return ParseResult(text="", error="Word 文档中没有可提取文本")
    return ParseResult(text=text, pages=[text])


def parse_pptx(data: bytes) -> ParseResult:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        pages.append(_clean("\n".join(texts)))
    text = _clean("\n\n".join(pages))
    if not text:
        return ParseResult(text="", error="PPT 中没有可提取文本")
    return ParseResult(text=text, pages=pages)


def parse_document(filename: str, data: bytes) -> ParseResult:
    name = filename.lower()
    if name.endswith(".pdf"):
        return parse_pdf(data)
    if name.endswith(".docx"):
        return parse_docx(data)
    if name.endswith(".pptx"):
        return parse_pptx(data)
    if name.endswith(".md"):
        return parse_markdown(data)
    if name.endswith(".txt"):
        text = data.decode("utf-8", errors="ignore")
        return ParseResult(text=_clean(text), pages=[_clean(text)])
    return ParseResult(text="", error="仅支持 PDF / Word / PPT / TXT / Markdown（.md）")
